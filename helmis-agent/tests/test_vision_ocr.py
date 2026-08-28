"""
test_vision_ocr.py — Unit tests for Multimodal Vision OCR and Hybrid Document Reader.
"""

import io
import os
import shutil
import tempfile
from collections.abc import Generator
from unittest.mock import MagicMock, patch

import pptx
import pymupdf
import pytest
from PIL import Image

from src.memory.ocr import async_perform_vision_ocr, perform_vision_ocr
from src.memory.vault import (
    init_vault_structure,
    read_vault_file,
    save_file_to_vault,
)


@pytest.fixture(autouse=True)
def isolated_vault_environment(monkeypatch: pytest.MonkeyPatch) -> Generator[str, None, None]:
    """Provide an isolated temporary directory for vault files and catalog during tests."""
    temp_dir = tempfile.mkdtemp(prefix="helmis_ocr_test_")
    vault_dir = os.path.join(temp_dir, "vault")
    catalog_file = os.path.join(temp_dir, "file_catalog.json")

    monkeypatch.setattr("src.memory.vault.DATA_DIR", temp_dir)
    monkeypatch.setattr("src.memory.vault.VAULT_DIR", vault_dir)
    monkeypatch.setattr("src.memory.vault.CATALOG_FILE", catalog_file)

    init_vault_structure()

    yield temp_dir

    shutil.rmtree(temp_dir, ignore_errors=True)


def test_perform_vision_ocr_sync_and_async() -> None:
    """Verify perform_vision_ocr and async_perform_vision_ocr parse Gemini JSON response."""
    mock_resp = MagicMock()
    mock_resp.status_code = 200
    mock_resp.json.return_value = {
        "candidates": [
            {
                "content": {
                    "parts": [{"text": "### Invoice 2026\n| Item | Harga |\n| VPS | 150.000 |"}]
                }
            }
        ]
    }

    dummy_img_bytes = b"fake_png_bytes"
    with patch("src.memory.ocr.get_next_gemini_key", return_value="AIzaFakeKey"):
        with patch("httpx.Client.post", return_value=mock_resp):
            res = perform_vision_ocr(dummy_img_bytes, "image/png")
            assert res is not None
            assert "### Invoice 2026" in res
            assert "150.000" in res

    # Test empty bytes returns None
    assert perform_vision_ocr(b"") is None


@pytest.mark.asyncio
async def test_async_perform_vision_ocr() -> None:
    """Verify async_perform_vision_ocr works with mock."""
    mock_resp = MagicMock()
    mock_resp.status_code = 200
    mock_resp.json.return_value = {
        "candidates": [
            {
                "content": {
                    "parts": [{"text": "KTP Republik Indonesia\nNIK: 320123456789"}]
                }
            }
        ]
    }

    dummy_img_bytes = b"fake_png_bytes"
    with patch("src.memory.ocr.get_next_gemini_key", return_value="AIzaFakeKey"):
        with patch("httpx.AsyncClient.post", return_value=mock_resp):
            res = await async_perform_vision_ocr(dummy_img_bytes, "image/jpeg")
            assert res is not None
            assert "NIK: 320123456789" in res


def test_read_scanned_pdf_triggers_vision_ocr() -> None:
    """Verify scanned/raster image-only PDF automatically triggers Vision OCR."""
    # 1. Create a 1-page PDF where the page only contains a raster image (no text layer)
    img = Image.new("RGB", (300, 150), color=(255, 255, 255))
    img_bytes_io = io.BytesIO()
    img.save(img_bytes_io, format="PNG")
    img_bytes = img_bytes_io.getvalue()

    doc = pymupdf.open()
    page = doc.new_page(width=300, height=150)
    page.insert_image(page.rect, stream=img_bytes)
    scanned_pdf_bytes = doc.tobytes()
    doc.close()

    # Save to vault
    rec = save_file_to_vault(
        data=scanned_pdf_bytes,
        filename="sk_rektor_scan.pdf",
        category="documents",
        owner="Gilang",
    )

    # Mock Vision OCR returning structured extracted text
    mock_ocr_result = "SURAT KEPUTUSAN REKTOR\nNomor: 888/SK/2026\nMenetapkan: Beasiswa Penuh"
    with patch("src.memory.vault.perform_vision_ocr", return_value=mock_ocr_result) as mock_ocr:
        res = read_vault_file(rec["id"])

        assert res["status"] == "success"
        assert res["content_type"] == "pdf"
        mock_ocr.assert_called_once()
        assert "[Hasil Vision OCR]" in res["content"]
        assert "SURAT KEPUTUSAN REKTOR" in res["content"]
        assert "Nomor: 888/SK/2026" in res["content"]


def test_read_digital_pdf_bypasses_vision_ocr() -> None:
    """Verify digital text PDF uses fast text extraction without invoking Vision OCR."""
    # Create digital PDF with embedded text
    doc = pymupdf.open()
    page = doc.new_page(width=400, height=300)
    page.insert_text(
        (50, 50),
        "Laporan Kemajuan Proyek Helmis Autonomous AI Agent Tahun 2026 Semester Ganjil.",
    )
    digital_pdf_bytes = doc.tobytes()
    doc.close()

    rec = save_file_to_vault(
        data=digital_pdf_bytes,
        filename="laporan_kemajuan.pdf",
        category="documents",
        owner="Gilang",
    )

    with patch("src.memory.vault.perform_vision_ocr") as mock_ocr:
        res = read_vault_file(rec["id"])
        assert res["status"] == "success"
        mock_ocr.assert_not_called()
        assert "Laporan Kemajuan Proyek Helmis" in res["content"]


def test_read_pptx_picture_slide_triggers_vision_ocr() -> None:
    """Verify PPTX slide with picture/diagram triggers Vision OCR when slide text is sparse."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    img = Image.new("RGB", (200, 100), color=(100, 150, 200))
    img_io = io.BytesIO()
    img.save(img_io, format="PNG")
    img_io.seek(0)
    slide.shapes.add_picture(img_io, 50000, 50000)

    buf = io.BytesIO()
    prs.save(buf)

    rec = save_file_to_vault(
        data=buf.getvalue(),
        filename="arsitektur_diagram.pptx",
        category="projects",
        owner="Gilang",
    )

    mock_diagram_ocr = "Diagram Alur: [Client WAHA] -> [Queue FIFO] -> [ReAct Engine]"
    with patch("src.memory.vault.perform_vision_ocr", return_value=mock_diagram_ocr) as mock_ocr:
        res = read_vault_file(rec["id"])
        assert res["status"] == "success"
        assert res["content_type"] == "pptx"
        mock_ocr.assert_called_once()
        assert "*(Hasil Vision OCR Gambar Slide)*" in res["content"]
        assert "Diagram Alur:" in res["content"]


def test_read_standalone_image_triggers_and_caches_ocr() -> None:
    """Verify reading a PNG image triggers Vision OCR and saves into catalog ocr_summary."""
    img = Image.new("RGB", (300, 100), color=(240, 240, 240))
    img_io = io.BytesIO()
    img.save(img_io, format="PNG")

    rec = save_file_to_vault(
        data=img_io.getvalue(),
        filename="struk_kopi.png",
        category="receipts",
        owner="Bunga",
    )

    mock_receipt_ocr = "Kopi Kenangan\n1x Americano: Rp 22.000\nTotal: Rp 22.000"
    with patch("src.memory.vault.perform_vision_ocr", return_value=mock_receipt_ocr) as mock_ocr:
        res = read_vault_file(rec["id"])
        assert res["status"] == "success"
        assert res["content_type"] == "image"
        mock_ocr.assert_called_once()
        assert "Kopi Kenangan" in res["content"]
        assert "Rp 22.000" in res["content"]

    # Subsequent read should use cached OCR without calling perform_vision_ocr again
    with patch("src.memory.vault.perform_vision_ocr") as mock_ocr_again:
        res2 = read_vault_file(rec["id"])
        mock_ocr_again.assert_not_called()
        assert "Kopi Kenangan" in res2["content"]

"""
vault.py — Core Document Vault, File Storage Hierarchy, Metadata Catalog, and Safe Dynamic Operations.
"""

import fcntl
import hashlib
import json
import logging
import mimetypes
import os
import re
import shutil
import time
import uuid
from datetime import datetime
from typing import Any, cast
from zoneinfo import ZoneInfo

from .ocr import perform_vision_ocr

log = logging.getLogger("helmis-vault")
TZ = ZoneInfo("Asia/Jakarta")

DATA_DIR = os.environ.get("DATA_DIR") or (
    "/app/data" if os.path.exists("/app/data") else os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "data")
)
VAULT_DIR = os.path.join(DATA_DIR, "vault")
CATALOG_FILE = os.path.join(DATA_DIR, "file_catalog.json")

DEFAULT_CATEGORIES = ["health", "id_cards", "travel", "receipts", "documents", "media", "projects"]
_env_owner1 = os.environ.get("OWNER_NAME", "").strip().lower()
_env_owner2 = os.environ.get("PARTNER_NAME", "").strip().lower()
DEFAULT_OWNERS = [o for o in set(["gilang", "bunga", "shared", _env_owner1, _env_owner2]) if o]


def _get_vault_dir() -> str:
    import sys
    v_dir = globals().get("VAULT_DIR") or os.path.join(DATA_DIR, "vault")
    for mod_name in ("src.vault", "src.memory", "src.memory.vault"):
        if mod_name in sys.modules:
            mod = sys.modules[mod_name]
            if hasattr(mod, "VAULT_DIR"):
                v_dir = getattr(mod, "VAULT_DIR")
    return v_dir


def _get_catalog_file() -> str:
    import sys
    c_file = globals().get("CATALOG_FILE") or os.path.join(DATA_DIR, "file_catalog.json")
    for mod_name in ("src.vault", "src.memory", "src.memory.vault"):
        if mod_name in sys.modules:
            mod = sys.modules[mod_name]
            if hasattr(mod, "CATALOG_FILE"):
                c_file = getattr(mod, "CATALOG_FILE")
    return c_file


def is_safe_vault_path(target_path: str) -> bool:
    """Verify that target_path is strictly within VAULT_DIR (prevent path traversal)."""
    try:
        abs_vault = os.path.abspath(_get_vault_dir())
        abs_target = os.path.abspath(target_path)
        return os.path.commonpath([abs_vault, abs_target]) == abs_vault
    except Exception:
        return False


def init_vault_structure() -> None:
    """Initialize base directory structure under VAULT_DIR if not present."""
    v_dir = _get_vault_dir()
    c_file = _get_catalog_file()
    os.makedirs(v_dir, exist_ok=True)
    for cat in DEFAULT_CATEGORIES:
        for owner in DEFAULT_OWNERS:
            os.makedirs(os.path.join(v_dir, cat, owner), exist_ok=True)
    if not os.path.exists(c_file):
        _save_catalog({"files": [], "version": 1})


def sanitize_filename(filename: str) -> str:
    """Sanitize filename to alphanumeric, underscores, hyphens, and standard extension."""
    clean = os.path.basename(filename).strip()
    name, ext = os.path.splitext(clean)
    name = re.sub(r"[^\w\-\.]+", "_", name).strip("._")
    if not name:
        name = f"file_{int(time.time())}"
    ext = ext.lower()
    return f"{name}{ext}"


def _load_catalog() -> dict[str, Any]:
    init_vault_structure()
    c_file = _get_catalog_file()
    if not os.path.exists(c_file):
        return {"files": [], "version": 1}
    try:
        with open(c_file, encoding="utf-8") as f:
            fcntl.flock(f.fileno(), fcntl.LOCK_SH)
            try:
                return cast(dict[str, Any], json.load(f))
            finally:
                fcntl.flock(f.fileno(), fcntl.LOCK_UN)
    except Exception as err:
        log.error("Error reading file_catalog.json: %s", err)
        return {"files": [], "version": 1}


def _sanitize_surrogates(obj: Any) -> Any:
    """Recursively sanitize surrogate characters from strings to guarantee valid UTF-8 JSON serialization."""
    if isinstance(obj, str):
        return obj.encode("utf-8", errors="replace").decode("utf-8", errors="replace")
    if isinstance(obj, dict):
        return {_sanitize_surrogates(k): _sanitize_surrogates(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_sanitize_surrogates(v) for v in obj]
    return obj


def _save_catalog(data: dict[str, Any]) -> None:
    c_file = _get_catalog_file()
    parent = os.path.dirname(c_file)
    if parent:
        os.makedirs(parent, exist_ok=True)
    tmp_file = f"{c_file}.tmp.{uuid.uuid4().hex[:8]}"
    clean_data = _sanitize_surrogates(data)
    with open(tmp_file, "w", encoding="utf-8", errors="replace") as f:
        fcntl.flock(f.fileno(), fcntl.LOCK_EX)
        try:
            json.dump(clean_data, f, indent=2, ensure_ascii=False)
            f.flush()
            os.fsync(f.fileno())
        finally:
            fcntl.flock(f.fileno(), fcntl.LOCK_UN)
    os.replace(tmp_file, c_file)


def save_file_to_vault(
    data: bytes,
    filename: str,
    owner: str = "Gilang",
    category: str = "documents",
    subfolder: str = "",
    description: str = "",
    tags: list[str] | None = None,
    ocr_summary: str = "",
    allow_versioning: bool = True,
    original_filename: str | None = None,
) -> dict[str, Any]:
    """Save raw file bytes into the Document Vault with automatic catalog indexing."""
    init_vault_structure()
    clean_name = sanitize_filename(filename)
    orig_name = (original_filename or filename).strip()
    # Strip any directory path traversal from orig_name for display safety
    orig_name = os.path.basename(orig_name) or filename
    owner_clean = owner.strip().lower()
    if owner_clean not in ("gilang", "bunga", "both", "shared"):
        owner_clean = "gilang"
    owner_folder = "shared" if owner_clean in ("both", "shared") else owner_clean
    category_clean = sanitize_filename(category).lower().strip("._") or "documents"

    content_hash = hashlib.sha256(data).hexdigest()

    catalog = _load_catalog()
    files = catalog.get("files", [])

    for existing in files:
        if existing.get("content_hash") == content_hash:
            log.info("Identical file already exists in vault: %s", existing.get("filename"))
            if orig_name and existing.get("original_filename") != orig_name:
                existing["original_filename"] = orig_name
                _save_catalog(catalog)
            return cast(dict[str, Any], existing)

    v_dir = _get_vault_dir()
    if subfolder:
        clean_subfolder = os.path.normpath(subfolder).strip("/\\")
        dest_dir = os.path.join(v_dir, clean_subfolder)
    else:
        dest_dir = os.path.join(v_dir, category_clean, owner_folder)

    if not is_safe_vault_path(dest_dir):
        dest_dir = os.path.join(v_dir, "documents", owner_folder)

    os.makedirs(dest_dir, exist_ok=True)

    final_name = clean_name
    dest_file = os.path.join(dest_dir, final_name)
    if os.path.exists(dest_file) and allow_versioning:
        base_stem, ext = os.path.splitext(clean_name)
        ver = 2
        while os.path.exists(os.path.join(dest_dir, f"{base_stem}_v{ver}{ext}")):
            ver += 1
        final_name = f"{base_stem}_v{ver}{ext}"
        dest_file = os.path.join(dest_dir, final_name)

    tmp_dest = f"{dest_file}.tmp.{uuid.uuid4().hex[:6]}"
    with open(tmp_dest, "wb") as f:
        f.write(data)
        f.flush()
        os.fsync(f.fileno())
    os.replace(tmp_dest, dest_file)

    mime_type, _ = mimetypes.guess_type(dest_file)
    if not mime_type:
        mime_type = "application/octet-stream"

    rel_path = os.path.relpath(dest_file, v_dir).replace("\\", "/")
    now_dt = datetime.now(TZ)
    now_str = now_dt.strftime("%A, %d %B %Y - %H:%M WIB")

    file_record: dict[str, Any] = {
        "id": f"doc_{int(time.time())}_{uuid.uuid4().hex[:6]}",
        "filename": final_name,
        "original_filename": orig_name,
        "relative_path": rel_path,
        "category": category_clean,
        "owner": "Both" if owner_folder == "shared" else owner.capitalize(),
        "mime_type": mime_type,
        "size_bytes": len(data),
        "content_hash": content_hash,
        "tags": tags or [],
        "description": description.strip() or orig_name,
        "ocr_summary": ocr_summary.strip(),
        "created_at": now_str,
        "updated_at": now_str,
    }

    files.append(file_record)
    catalog["files"] = files
    _save_catalog(catalog)
    log.info("Saved file to vault: %s (%s bytes)", rel_path, len(data))
    return file_record


def search_vault(
    query: str,
    owner: str | None = None,
    category: str | None = None,
    limit: int = 10,
) -> list[dict[str, Any]]:
    """Search files in vault by keyword across filename, original_filename, description, tags, and OCR summary."""
    catalog = _load_catalog()
    files = catalog.get("files", [])
    q = query.lower().strip()
    owner_filter = owner.strip().lower() if owner else None
    cat_filter = category.strip().lower() if category else None

    matches: list[dict[str, Any]] = []
    for f in files:
        if owner_filter:
            f_owner = str(f.get("owner", "")).lower()
            if owner_filter not in f_owner and f_owner not in ("both", "shared"):
                continue
        if cat_filter:
            f_cat = str(f.get("category", "")).lower()
            if cat_filter != f_cat and cat_filter not in str(f.get("relative_path", "")).lower():
                continue

        fn = str(f.get("filename", "")).lower()
        orig_fn = str(f.get("original_filename", "")).lower()
        desc = str(f.get("description", "")).lower()
        tags = [str(t).lower() for t in f.get("tags", [])]
        ocr = str(f.get("ocr_summary", "")).lower()

        searchable_text = f"{fn} {orig_fn} {desc} {' '.join(tags)} {ocr}"
        if q in fn or q in orig_fn or q in desc or any(q in t for t in tags) or q in ocr:
            matches.append(f)
        elif all(word in searchable_text for word in q.split()):
            matches.append(f)

    return matches[:limit]


def list_vault_files(
    owner: str | None = None,
    category: str | None = None,
    directory: str | None = None,
) -> list[dict[str, Any]]:
    """List files registered in the Document Vault filtered by owner, category, or subfolder."""
    catalog = _load_catalog()
    files = catalog.get("files", [])
    owner_filter = owner.strip().lower() if owner else None
    cat_filter = category.strip().lower() if category else None
    dir_filter = directory.strip().lower().replace("\\", "/") if directory else None

    result: list[dict[str, Any]] = []
    for f in files:
        if owner_filter:
            f_owner = str(f.get("owner", "")).lower()
            if owner_filter not in f_owner and f_owner not in ("both", "shared"):
                continue
        if cat_filter:
            f_cat = str(f.get("category", "")).lower()
            if cat_filter != f_cat and cat_filter not in str(f.get("relative_path", "")).lower():
                continue
        if dir_filter:
            rel = str(f.get("relative_path", "")).lower()
            if not rel.startswith(dir_filter):
                continue
        result.append(f)
    return result


def get_vault_file_by_id(file_id: str) -> tuple[dict[str, Any], bytes] | None:
    """Retrieve file record and raw bytes by file_id."""
    catalog = _load_catalog()
    v_dir = _get_vault_dir()
    for f in catalog.get("files", []):
        if f.get("id") == file_id:
            full_path = os.path.join(v_dir, f.get("relative_path", ""))
            if os.path.exists(full_path) and is_safe_vault_path(full_path):
                with open(full_path, "rb") as fp:
                    return f, fp.read()
    return None


def get_vault_file_by_name(filename: str, owner: str | None = None) -> tuple[dict[str, Any], bytes] | None:
    """Retrieve file record and raw bytes by filename match."""
    matches = search_vault(query=filename, owner=owner, limit=1)
    if matches:
        return get_vault_file_by_id(matches[0]["id"])
    return None


def _resolve_target_files(target: str | list[str], exact_only: bool = False) -> list[dict[str, Any]]:
    """Helper to resolve a target (ID, filename, list of IDs, or search query) to file records."""
    catalog = _load_catalog()
    files = catalog.get("files", [])
    if isinstance(target, list):
        ids_or_names = {t.lower().strip() for t in target}
        return [f for f in files if f.get("id", "").lower() in ids_or_names or f.get("filename", "").lower() in ids_or_names]

    t_str = target.strip()
    if not t_str:
        return []

    # 1. Exact ID match
    for f in files:
        if f.get("id") == t_str:
            return [f]

    # 2. Exact filename match
    for f in files:
        if f.get("filename", "").lower() == t_str.lower():
            return [f]

    if exact_only:
        return []

    # 3. Fallback to keyword search across catalog (requires min 3 chars to prevent accidental broad wipes)
    if len(t_str) >= 3:
        return search_vault(query=t_str, limit=50)
    return []


def move_vault_files(
    target: str | list[str],
    destination_dir: str | None = None,
    new_category: str | None = None,
    new_owner: str | None = None,
) -> list[dict[str, Any]]:
    """
    Polymorphic Move: Moves 1 or many files to a destination directory, category, or owner.
    Target can be a single file ID/name, a list of IDs/names, or a search query string.
    """
    catalog = _load_catalog()
    files = catalog.get("files", [])
    targets = _resolve_target_files(target)

    if not targets:
        return []

    moved_records: list[dict[str, Any]] = []
    target_ids = {t["id"] for t in targets}
    v_dir = _get_vault_dir()

    for record in files:
        if record["id"] not in target_ids:
            continue

        old_rel = record.get("relative_path", "")
        old_full = os.path.join(v_dir, old_rel)
        if not os.path.exists(old_full) or not is_safe_vault_path(old_full):
            continue

        cat = (new_category or record.get("category", "documents")).lower().strip()
        owner_str = (new_owner or record.get("owner", "Gilang")).capitalize()
        owner_folder = "shared" if owner_str.lower() in ("both", "shared") else owner_str.lower()

        if destination_dir:
            dest_dir = os.path.join(v_dir, destination_dir.strip("/\\"))
        else:
            dest_dir = os.path.join(v_dir, cat, owner_folder)

        if not is_safe_vault_path(dest_dir):
            continue

        os.makedirs(dest_dir, exist_ok=True)
        filename = record.get("filename", "")
        new_full = os.path.join(dest_dir, filename)

        # Collision versioning if destination already exists and is different file
        if os.path.exists(new_full) and os.path.abspath(old_full) != os.path.abspath(new_full):
            base_stem, ext = os.path.splitext(filename)
            ver = 2
            while os.path.exists(os.path.join(dest_dir, f"{base_stem}_v{ver}{ext}")):
                ver += 1
            filename = f"{base_stem}_v{ver}{ext}"
            new_full = os.path.join(dest_dir, filename)

        shutil.move(old_full, new_full)
        new_rel = os.path.relpath(new_full, v_dir).replace("\\", "/")

        record["filename"] = filename
        record["relative_path"] = new_rel
        record["category"] = cat
        record["owner"] = owner_str
        record["updated_at"] = datetime.now(TZ).strftime("%A, %d %B %Y - %H:%M WIB")
        moved_records.append(record)
        log.info("Moved vault file from %s to %s", old_rel, new_rel)

    if moved_records:
        _save_catalog(catalog)

    return moved_records


def delete_vault_files(target: str | list[str], allow_bulk_query: bool = True) -> list[str]:
    """
    Polymorphic Delete: Deletes 1 or many files matched by ID, filename, list, or search query.
    Returns list of deleted filenames.
    """
    catalog = _load_catalog()
    files = catalog.get("files", [])
    targets = _resolve_target_files(target, exact_only=not allow_bulk_query)

    if not targets:
        return []

    target_ids = {t["id"] for t in targets}
    deleted_names: list[str] = []
    v_dir = _get_vault_dir()

    remaining_files: list[dict[str, Any]] = []
    for f in files:
        if f.get("id") in target_ids:
            full_path = os.path.join(v_dir, f.get("relative_path", ""))
            if os.path.exists(full_path) and is_safe_vault_path(full_path):
                try:
                    os.remove(full_path)
                except Exception as err:
                    log.error("Failed to remove file from disk: %s", err)
            deleted_names.append(f.get("filename", ""))
        else:
            remaining_files.append(f)

    if deleted_names:
        catalog["files"] = remaining_files
        _save_catalog(catalog)

    return deleted_names


def create_vault_directory(dir_path: str) -> str:
    """Create a new custom directory path under the Document Vault."""
    init_vault_structure()
    clean_path = dir_path.strip().replace("..", "").strip("/\\")
    v_dir = _get_vault_dir()
    full_path = os.path.join(v_dir, clean_path)
    if not is_safe_vault_path(full_path):
        raise ValueError(f"Invalid directory path: {dir_path}")
    os.makedirs(full_path, exist_ok=True)
    return os.path.relpath(full_path, v_dir).replace("\\", "/")


def delete_vault_directory(dir_path: str, recursive: bool = False) -> tuple[bool, str]:
    """Delete a directory from the vault (empty or recursive with catalog cleanup)."""
    init_vault_structure()
    clean_path = dir_path.strip().replace("..", "").strip("/\\")
    v_dir = _get_vault_dir()
    full_path = os.path.join(v_dir, clean_path)
    if not is_safe_vault_path(full_path):
        return False, "Direktori di luar brankas tidak dapat dihapus."

    if not os.path.exists(full_path):
        return False, f"Direktori '{clean_path}' tidak ditemukan."

    if full_path == os.path.abspath(v_dir):
        return False, "Direktori root brankas tidak dapat dihapus."

    # Guard core categories from deletion
    if clean_path.lower() in DEFAULT_CATEGORIES:
        return False, f"Kategori utama '{clean_path}' dilindungi dan tidak dapat dihapus."

    entries = os.listdir(full_path)
    if entries and not recursive:
        return False, f"Direktori '{clean_path}' tidak kosong ({len(entries)} item). Gunakan recursive=True untuk menghapus beserta isinya."

    catalog = _load_catalog()
    files = catalog.get("files", [])
    rel_prefix = os.path.relpath(full_path, v_dir).replace("\\", "/")

    updated_files = [f for f in files if not str(f.get("relative_path", "")).startswith(rel_prefix)]
    removed_count = len(files) - len(updated_files)

    shutil.rmtree(full_path)
    catalog["files"] = updated_files
    _save_catalog(catalog)

    return True, f"Direktori '{clean_path}' berhasil dihapus ({removed_count} file di-unregister)."


def _update_vault_file_ocr(file_id: str, ocr_text: str) -> None:
    """Update OCR summary metadata for a file in catalog."""
    try:
        catalog = _load_catalog()
        for f in catalog.get("files", []):
            if f.get("id") == file_id:
                f["ocr_summary"] = ocr_text[:2000]
                _save_catalog(catalog)
                break
    except Exception as ex:
        log.warning("Failed to update vault file OCR for %s: %s", file_id, ex)


def read_vault_file(
    file_id_or_name: str,
    max_chars: int = 15000,
    force_ocr: bool = False,
) -> dict[str, Any]:
    """
    Read content of a file from the vault.
    - Plain text / JSON / Markdown / CSV: Decoded as string.
    - PDF: Page-by-page digital text layer, falling back to Gemini Vision OCR for raster scans,
           or forcing full visual Vision OCR on all pages when force_ocr=True.
    - DOCX / PPTX / XLSX: Extracted to structured Markdown.
    - Images: Extracted via dynamic Gemini Vision OCR and cached into catalog.
    """
    res = get_vault_file_by_id(file_id_or_name)
    if not res:
        res = get_vault_file_by_name(file_id_or_name)
    if not res:
        return {"status": "error", "error": f"File '{file_id_or_name}' tidak ditemukan di brankas dokumen."}

    record, raw_bytes = res
    filename = record.get("filename", "")
    mime = record.get("mime_type", "")
    ext = os.path.splitext(filename)[1].lower()

    content_type = "binary"
    extracted_text = ""

    # 1. Plain text / Markdown / JSON / CSV / Code / Configs
    if (
        mime.startswith("text/")
        or ext in (
            ".txt",
            ".md",
            ".json",
            ".csv",
            ".tsv",
            ".py",
            ".yaml",
            ".yml",
            ".env",
            ".log",
            ".html",
            ".css",
            ".js",
        )
    ):
        content_type = "text"
        extracted_text = raw_bytes.decode("utf-8", errors="ignore").strip()

        # Binary spoof detection: if text contains null bytes or binary characters, treat as binary
        if raw_bytes and (b"\x00" in raw_bytes[:1024] or (extracted_text.count("\x00") / max(1, len(extracted_text)) > 0.02)):
            content_type = "binary"
            desc = record.get("description", "")
            extracted_text = f"[File Biner {mime}]: {desc} (Data biner mentah tidak dapat ditampilkan sebagai teks)"

    # 2. PDF Documents (.pdf) — PyMuPDF with High-Precision Fallback & Force Vision OCR
    elif ext == ".pdf" or mime == "application/pdf":
        content_type = "pdf"
        scanned_page_count = 0
        pdf_pages_text: list[str] = []
        is_encrypted = False
        try:
            import pymupdf

            doc = pymupdf.open(stream=raw_bytes, filetype="pdf")
            if doc.is_encrypted:
                try:
                    doc.authenticate("")
                except Exception:
                    is_encrypted = True

            if is_encrypted:
                extracted_text = "[Dokumen PDF terenkripsi / ber-password. Konten teks tidak dapat diekstrak tanpa kata sandi.]"
            else:
                total_p = len(doc)
                for i, page in enumerate(doc):
                    raw_txt = page.get_text("text")
                    p_txt = str(raw_txt).strip() if raw_txt else ""
                    page_parts: list[str] = []

                    # If force_ocr is enabled, bypass text layer and process whole page as image
                    if force_ocr:
                        scanned_page_count += 1
                        try:
                            pix = page.get_pixmap(dpi=150)
                            rendered_png = pix.tobytes("png")
                            ocr_txt = perform_vision_ocr(
                                rendered_png,
                                "image/png",
                                prompt_hint="Extract all readable text, timelines, schedules, tables, deadlines, milestones, and visual details from this document page image verbatim into clean Markdown.",
                            )
                            if ocr_txt and ocr_txt.strip():
                                page_parts.append(f"[Hasil Vision OCR (Image Mode)]\n{ocr_txt.strip()}")
                            elif p_txt:
                                page_parts.append(p_txt)
                        except Exception as ocr_err:
                            log.warning("Force Vision OCR failed on PDF page %d of %s: %s", i + 1, filename, ocr_err)
                            if p_txt:
                                page_parts.append(p_txt)
                    else:
                        if p_txt:
                            page_parts.append(p_txt)

                        # Inspect embedded images on this page
                        try:
                            img_list = page.get_images(full=True)
                        except Exception:
                            img_list = []

                        # If page has digital text AND embedded meaningful diagram/chart image(s):
                        if p_txt and img_list:
                            processed_imgs = 0
                            for img_info in img_list:
                                if processed_imgs >= 2:
                                    break
                                xref = img_info[0]
                                try:
                                    base_img = doc.extract_image(xref)
                                    width = base_img.get("width", 0)
                                    height = base_img.get("height", 0)
                                    img_data = base_img.get("image", b"")
                                    img_ext = base_img.get("ext", "png")
                                    # Filter out tiny decorative icons (only process meaningful diagrams >= 50x50 and >= 100 bytes)
                                    if width >= 50 and height >= 50 and len(img_data) >= 100:
                                        img_mime = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"
                                        ocr_sub = perform_vision_ocr(
                                            img_data,
                                            img_mime,
                                            prompt_hint="Extract all readable text, mathematical formulas ($...$), graphs, labels, numbers, flowchart nodes, and data from this diagram/chart image on the page into clean Markdown.",
                                        )
                                        if ocr_sub and ocr_sub.strip():
                                            page_parts.append(f"*(Hasil Vision OCR Diagram/Gambar Halaman {i+1})*:\n{ocr_sub.strip()}")
                                            scanned_page_count += 1
                                            processed_imgs += 1
                                except Exception as img_ex:
                                    log.warning("Failed to extract embedded image %d on page %d: %s", xref, i + 1, img_ex)

                        # If page has NO digital text (raster scan), render whole page:
                        elif not p_txt:
                            scanned_page_count += 1
                            try:
                                pix = page.get_pixmap(dpi=150)
                                rendered_png = pix.tobytes("png")
                                ocr_txt = perform_vision_ocr(
                                    rendered_png,
                                    "image/png",
                                    prompt_hint="Extract all readable text, tabular data, forms, signatures, and stamps from this scanned document page image into clean Markdown.",
                                )
                                if ocr_txt and ocr_txt.strip():
                                    page_parts.append(f"[Hasil Vision OCR]\n{ocr_txt.strip()}")
                            except Exception as ocr_err:
                                log.warning("Vision OCR failed on PDF page %d of %s: %s", i + 1, filename, ocr_err)

                    if page_parts:
                        pdf_pages_text.append(f"--- Halaman {i+1} dari {total_p} ---\n" + "\n\n".join(page_parts))

                doc.close()
                if pdf_pages_text:
                    extracted_text = "\n\n".join(pdf_pages_text)
        except Exception as ex:
            log.warning("pymupdf text extraction error on %s: %s, falling back to pypdf", filename, ex)
            try:
                import io
                import pypdf

                reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
                total_p = len(reader.pages)
                for i, page in enumerate(reader.pages):
                    p_txt = page.extract_text() or ""
                    if p_txt.strip():
                        pdf_pages_text.append(f"--- Halaman {i+1} dari {total_p} ---\n{p_txt.strip()}")
                if pdf_pages_text:
                    extracted_text = "\n\n".join(pdf_pages_text)
            except Exception as fb_ex:
                log.warning("pypdf fallback error on %s: %s", filename, fb_ex)

        # Fallback to OCR summary if PDF has no embedded text layer
        if not extracted_text.strip():
            ocr = record.get("ocr_summary", "")
            if ocr:
                extracted_text = f"[Hasil OCR Dokumen Scan]:\n{ocr}"
            else:
                extracted_text = "[Dokumen PDF raster scan tanpa teks digital. Tidak ada teks yang dapat diekstrak.]"
        elif scanned_page_count > 0 and not record.get("ocr_summary"):
            _update_vault_file_ocr(record.get("id", ""), extracted_text[:1000])

    # 3. Microsoft Word Documents (.docx, .doc)
    elif ext in (".docx", ".doc") or "wordprocessingml" in mime:
        content_type = "docx"
        try:
            import io

            import docx

            doc = docx.Document(io.BytesIO(raw_bytes))
            doc_parts: list[str] = []
            for p in doc.paragraphs:
                p_txt = p.text.strip()
                if not p_txt:
                    continue
                if getattr(p, "style", None) and getattr(p.style, "name", "").startswith("Heading"):
                    doc_parts.append(f"\n### {p_txt}\n")
                else:
                    doc_parts.append(p_txt)

            for table in doc.tables:
                tbl_rows: list[str] = []
                for r_idx, row in enumerate(table.rows):
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    tbl_rows.append("| " + " | ".join(cells) + " |")
                    if r_idx == 0:
                        tbl_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                if tbl_rows:
                    doc_parts.append("\n" + "\n".join(tbl_rows) + "\n")

            extracted_text = "\n\n".join(doc_parts).strip()
            if not extracted_text:
                extracted_text = "[Dokumen Word kosong / tanpa teks terbaca]"
        except Exception as ex:
            log.warning("python-docx extraction error on %s: %s", filename, ex)
            desc = record.get("description", "")
            extracted_text = f"[Dokumen Word ({filename})]: {desc} (Gagal membaca teks dokumen: {ex})"

    # 4. Microsoft PowerPoint Presentations (.pptx, .ppt)
    elif ext in (".pptx", ".ppt") or "presentationml" in mime:
        content_type = "pptx"
        try:
            import io

            import pptx
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = pptx.Presentation(io.BytesIO(raw_bytes))
            slides_text: list[str] = []
            total_slides = len(prs.slides)
            for idx, slide in enumerate(prs.slides, start=1):
                slide_parts: list[str] = []
                if getattr(slide.shapes, "title", None) and slide.shapes.title.text:
                    slide_parts.append(f"**Judul:** {slide.shapes.title.text.strip()}")

                picture_shapes = []
                for shape in slide.shapes:
                    if shape == getattr(slide.shapes, "title", None):
                        continue
                    tf = getattr(shape, "text_frame", None)
                    tbl = getattr(shape, "table", None)
                    if tf and hasattr(tf, "paragraphs"):
                        for para in getattr(tf, "paragraphs", []):
                            ptxt = "".join(getattr(r, "text", "") for r in getattr(para, "runs", [])).strip()
                            if ptxt:
                                indent = "  " * (getattr(para, "level", 0) or 0)
                                slide_parts.append(f"{indent}- {ptxt}")
                    elif tbl and hasattr(tbl, "rows"):
                        tbl_rows = []
                        for r_idx, row in enumerate(getattr(tbl, "rows", [])):
                            cells = [getattr(c, "text", "").strip().replace("\n", " ") for c in getattr(row, "cells", [])]
                            tbl_rows.append("| " + " | ".join(cells) + " |")
                            if r_idx == 0:
                                tbl_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        if tbl_rows:
                            slide_parts.append("\n" + "\n".join(tbl_rows))
                    elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                        picture_shapes.append(shape)

                if getattr(slide, "has_notes_slide", False) and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    ntxt = slide.notes_slide.notes_text_frame.text.strip()
                    if ntxt:
                        slide_parts.append(f"*(Catatan Presenter: {ntxt})*")

                # If slide contains meaningful picture/diagram shapes, run Vision OCR on images!
                if picture_shapes:
                    processed_shapes = 0
                    for p_shape in picture_shapes:
                        if processed_shapes >= 2:
                            break
                        try:
                            img_blob = p_shape.image.blob
                            img_mime = getattr(p_shape.image, "content_type", None) or "image/png"
                            if len(img_blob) >= 100:  # Process valid image blobs
                                img_ocr = perform_vision_ocr(
                                    img_blob,
                                    img_mime,
                                    prompt_hint="Extract all readable text, mathematical formulas ($...$), graphs, flowchart nodes, diagram labels, and tables from this slide picture into clean Markdown.",
                                )
                                if img_ocr and img_ocr.strip():
                                    slide_parts.append(f"*(Hasil Vision OCR Gambar/Diagram Slide)*:\n{img_ocr.strip()}")
                                    processed_shapes += 1
                        except Exception as img_err:
                            log.warning("Vision OCR failed on slide picture %d of %s: %s", idx, filename, img_err)

                if slide_parts:
                    slides_text.append(f"--- Slide {idx} dari {total_slides} ---\n" + "\n".join(slide_parts))
                else:
                    slides_text.append(f"--- Slide {idx} dari {total_slides} ---\n*(Slide tanpa teks/hanya gambar)*")

            extracted_text = "\n\n".join(slides_text).strip()
            if not extracted_text:
                extracted_text = "[Presentasi PowerPoint kosong / tanpa teks terbaca]"
        except Exception as ex:
            log.warning("python-pptx extraction error on %s: %s", filename, ex)
            desc = record.get("description", "")
            extracted_text = f"[Presentasi PowerPoint ({filename})]: {desc} (Gagal membaca teks slide: {ex})"

    # 5. Microsoft Excel Spreadsheets (.xlsx, .xls)
    elif ext in (".xlsx", ".xls") or "spreadsheetml" in mime:
        content_type = "xlsx"
        try:
            import io

            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True, read_only=True)
            sheets_text: list[str] = []
            for sname in wb.sheetnames:
                ws = wb[sname]
                rows_data: list[list[str]] = []
                for r in ws.iter_rows(values_only=True):
                    if any(c is not None and str(c).strip() != "" for c in r):
                        rows_data.append([str(c) if c is not None else "" for c in r])
                if not rows_data:
                    sheets_text.append(f"### Sheet: {sname}\n*(Sheet kosong)*")
                    continue
                capped_rows = rows_data[:100]
                headers = [h.strip().replace("\n", " ") for h in capped_rows[0]]
                tbl_lines = [
                    "| " + " | ".join(headers) + " |",
                    "| " + " | ".join(["---"] * len(headers)) + " |",
                ]
                for row in capped_rows[1:]:
                    tbl_lines.append("| " + " | ".join(c.strip().replace("\n", " ") for c in row) + " |")
                extra = f"\n*(Menampilkan {len(capped_rows)} dari {len(rows_data)} baris total)*" if len(rows_data) > 100 else ""
                sheets_text.append(f"### Sheet: {sname}\n" + "\n".join(tbl_lines) + extra)
            wb.close()
            extracted_text = "\n\n".join(sheets_text).strip()
            if not extracted_text:
                extracted_text = "[Spreadsheet Excel kosong / tanpa tabel terbaca]"
        except Exception as ex:
            log.warning("openpyxl extraction error on %s: %s", filename, ex)
            desc = record.get("description", "")
            extracted_text = f"[Spreadsheet Excel ({filename})]: {desc} (Gagal membaca tabel sheet: {ex})"

    # 6. Images / Other media
    elif mime.startswith("image/"):
        content_type = "image"
        ocr = record.get("ocr_summary", "")
        desc = record.get("description", "")

        # Dynamic Vision OCR extraction if summary is missing
        if not ocr or ocr.strip() in ("", "Tidak ada ringkasan OCR."):
            try:
                fresh_ocr = perform_vision_ocr(
                    raw_bytes,
                    mime,
                    prompt_hint="Extract all text, receipt items, prices, dates, tabular numbers, and forms from this document image into clean Markdown.",
                )
                if fresh_ocr and fresh_ocr.strip():
                    ocr = fresh_ocr.strip()
                    record["ocr_summary"] = ocr
                    _update_vault_file_ocr(record.get("id", ""), ocr)
            except Exception as ocr_e:
                log.warning("Vision OCR failed on image %s: %s", filename, ocr_e)

        extracted_text = f"[File Gambar/Foto]: {desc}\n\n[Hasil Vision OCR]:\n{ocr if ocr else 'Tidak ada teks yang dapat diekstrak dari gambar ini.'}"

    else:
        content_type = "binary"
        desc = record.get("description", "")
        extracted_text = f"[File Biner {mime}]: {desc}"

    # Apply length clipping if exceeded
    effective_max = max(100, max_chars)
    is_truncated = False
    if len(extracted_text) > effective_max:
        extracted_text = (
            extracted_text[:effective_max]
            + f"\n\n... (Dipotong karena melebihi batas {effective_max} karakter)"
        )
        is_truncated = True

    return {
        "status": "success",
        "file": record,
        "content_type": content_type,
        "content": extracted_text,
        "is_truncated": is_truncated,
        "size_bytes": len(raw_bytes),
    }


def list_vault_directories() -> list[str]:
    """List all directory paths inside the Document Vault."""
    init_vault_structure()
    dirs: list[str] = []
    v_dir = _get_vault_dir()
    for root, dirnames, _ in os.walk(v_dir):
        for d in dirnames:
            full = os.path.join(root, d)
            rel = os.path.relpath(full, v_dir).replace("\\", "/")
            dirs.append(rel)
    return sorted(dirs)
